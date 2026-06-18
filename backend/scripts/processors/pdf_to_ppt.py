import os
from pathlib import Path
from .utils import pptx, extract_pdf_text

# python-pptx EMU conversions
EMU_PER_PT = 12700
EMU_PER_INCH = 914400

def _pt(points): return int(points * EMU_PER_PT)
def _in(inches): return int(inches * EMU_PER_INCH)


def pdf_to_ppt(input_paths, output_path):
    if not input_paths:
        return
    pdf_text = extract_pdf_text(input_paths[0])
    layout = os.environ.get("PDF_TO_PPT_LAYOUT", "auto").lower()
    borders = os.environ.get("PDF_TO_PPT_BORDERS", "true").lower() == "true"

    if pptx:
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

        prs = pptx.Presentation()
        # 16:9 slide size
        prs.slide_width = _in(10)
        prs.slide_height = _in(5.63)

        blank_layout = prs.slide_layouts[6]  # blank layout

        # Split content into paragraphs (non-empty lines)
        lines = [ln.strip() for ln in pdf_text.split("\n") if ln.strip()]

        # Group 6 lines per slide
        lines_per_slide = 6
        slide_groups = [lines[i:i+lines_per_slide] for i in range(0, max(len(lines), 1), lines_per_slide)]

        # Color theme cycling
        bg_colors = [
            RGBColor(0x1e, 0x2d, 0x3d),
            RGBColor(0x2c, 0x3e, 0x50),
            RGBColor(0x1a, 0x1a, 0x2e),
        ]
        accent_colors = [
            RGBColor(0xe7, 0x4c, 0x3c),
            RGBColor(0x34, 0x98, 0xdb),
            RGBColor(0x2e, 0xcc, 0x71),
        ]

        for slide_idx, group in enumerate(slide_groups):
            slide = prs.slides.add_slide(blank_layout)

            bg_color = bg_colors[slide_idx % len(bg_colors)]
            accent = accent_colors[slide_idx % len(accent_colors)]

            # Slide background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

            # Accent top bar
            top_bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                0, 0, prs.slide_width, _pt(6)
            )
            try:
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = accent
                top_bar.line.fill.background()
            except Exception:
                pass

            # Slide number (bottom right)
            num_box = slide.shapes.add_textbox(
                prs.slide_width - _in(0.8), prs.slide_height - _pt(24),
                _in(0.7), _pt(20)
            )
            nf = num_box.text_frame
            nr = nf.paragraphs[0].add_run()
            nr.text = str(slide_idx + 1)
            nr.font.size = Pt(9)
            nr.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
            nr.font.bold = False
            nf.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # Title: first line of the group (or slide number)
            title_text = group[0] if group else f"Page {slide_idx + 1}"
            body_lines = group[1:] if len(group) > 1 else []

            title_box = slide.shapes.add_textbox(
                _in(0.4), _pt(20), prs.slide_width - _in(0.8), _pt(48)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tp = tf.paragraphs[0]
            tr = tp.add_run()
            tr.text = title_text[:120]
            tr.font.bold = True
            tr.font.size = Pt(22)
            tr.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

            # Body text box
            if body_lines:
                body_box = slide.shapes.add_textbox(
                    _in(0.4), _pt(90), prs.slide_width - _in(0.8), prs.slide_height - _pt(130)
                )
                bf = body_box.text_frame
                bf.word_wrap = True
                for i, line in enumerate(body_lines):
                    if i == 0:
                        bp = bf.paragraphs[0]
                    else:
                        bp = bf.add_paragraph()
                    br = bp.add_run()
                    br.text = "• " + line[:150]
                    br.font.size = Pt(13)
                    br.font.color.rgb = RGBColor(0xcc, 0xdd, 0xee)
                    bp.space_before = Pt(4)

            # Accent separator line under title (using a thin rectangle)
            try:
                sep = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    _in(0.4), _pt(82),
                    _in(3), _pt(2)
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = accent
                sep.line.fill.background()
            except Exception:
                pass

        prs.save(output_path)
    else:
        # Plain text fallback
        output_path.write_text(
            f"PDF → PowerPoint Export\nLayout: {layout.upper()}\nBorders: {borders}\n\n" + pdf_text,
            encoding="utf-8"
        )

    print(f"Converted PDF to Slide deck: {output_path}")
